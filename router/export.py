import io
import traceback
from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Optional
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from fpdf import FPDF

router = APIRouter(prefix="/api/v1/export", tags=["Export Tools"])

# --- PPTX ENDPOINT ---
class SlideData(BaseModel):
    slide_number: Optional[int] = None
    title: str
    bullet_points: List[str]
    speaker_notes: Optional[str] = ""

class PresentationExportPayload(BaseModel):
    presentation_title: str
    slides: List[SlideData]

@router.post("/pptx")
async def export_to_pptx(payload: PresentationExportPayload):
    try:
        prs = Presentation()
        
        # --- Theme Colors ---
        BG_COLOR = RGBColor(15, 23, 42)       # Deep Slate/Navy Background
        TITLE_COLOR = RGBColor(245, 158, 11)  # Amber/Gold Titles
        TEXT_COLOR = RGBColor(241, 245, 249)  # Off-White Body Text
        SUB_COLOR = RGBColor(148, 163, 184)   # Gray Subtitles
        
        # --- 1. Title Slide ---
        title_slide_layout = prs.slide_layouts[0] 
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Programmatic Background Color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Format Title
        title_shape = slide.shapes.title
        title_shape.text = payload.presentation_title
        title_font = title_shape.text_frame.paragraphs[0].font
        title_font.color.rgb = TITLE_COLOR
        title_font.bold = True
        
        # Format Subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = "Generated via Academic Architect AI"
        sub_font = subtitle_shape.text_frame.paragraphs[0].font
        sub_font.color.rgb = SUB_COLOR
        sub_font.italic = True

        # --- 2. Content Slides ---
        bullet_slide_layout = prs.slide_layouts[1]
        for slide_data in payload.slides:
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Programmatic Background Color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = BG_COLOR
            
            # Format Title
            title_shape = slide.shapes.title
            title_shape.text = slide_data.title
            title_font = title_shape.text_frame.paragraphs[0].font
            title_font.color.rgb = TITLE_COLOR
            title_font.bold = True
            
            # Format Bullet Points
            tf = slide.placeholders[1].text_frame
            tf.clear() # Clear default empty paragraph
            
            for point in slide_data.bullet_points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.color.rgb = TEXT_COLOR
                p.font.size = Pt(18)
                
            if slide_data.speaker_notes:
                slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)

        return StreamingResponse(
            stream, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={payload.presentation_title.replace(' ', '_')}.pptx"}
        )
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PPTX Generation failed: {str(e)}")


# --- PDF ENDPOINT ---

def safe_text(text):
    if text is None: return ""
    text = str(text)
    replacements = {
        '\u2018': "'", '\u2019': "'", '\u201c': '"', '\u201d': '"', 
        '\u2013': '-', '\u2014': '-', '\u2026': '...', '\u00a0': ' ', 
        '\u03c0': 'pi', '\u00b0': ' degrees', '\u00d7': 'x', 
        '\u2220': 'Angle ', '\u2264': '<=', '\u2265': '>=', '\u2260': '!=', 
        '\u2192': '->', '\u221a': 'sqrt',
        '\u22a5': ' perpendicular to ', '\u25b3': 'Triangle ', '\u2206': 'Triangle ', '\u0394': 'Triangle ',
        '\u2245': '~=', '\u223c': '~', '\u2234': 'Therefore ', '\u2235': 'Because ',
        '•': '-', '\u2022': '-', '\u25CF': '-', '\u25CB': '-', 
        '\u00B7': '-', '\u25AA': '-', '\u2043': '-', '\u2023': '-',
        '\u27c2': ' perpendicular to ', '⊥': ' perpendicular to '
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text.encode('latin-1', 'replace').decode('latin-1')

class PDFReport(FPDF):
    def header(self):
        self.set_font('helvetica', 'B', 16)
        self.set_text_color(0, 88, 190)
        self.cell(0, 10, 'Academic Architect AI Export', align='C')
        self.ln(15)

    def footer(self):
        self.set_y(-15)
        self.set_font('helvetica', 'I', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, f'Page {self.page_no()}', align='C')

# The Generic Formatter (Used by Lesson Plan, Study Notes, etc.)
def add_json_to_pdf(pdf, key, value):
    if isinstance(value, dict):
        if key:
            pdf.set_font("helvetica", "B", 14)
            pdf.set_text_color(11, 28, 48)
            pdf.set_x(10)
            pdf.multi_cell(0, 8, txt=safe_text(str(key).replace("_", " ").title()))
            pdf.ln(2)
        for k, v in value.items():
            add_json_to_pdf(pdf, k, v)
        pdf.ln(4)
    elif isinstance(value, list):
        if key:
            pdf.set_font("helvetica", "B", 12)
            pdf.set_text_color(11, 28, 48)
            pdf.set_x(10) 
            pdf.multi_cell(0, 8, txt=safe_text(str(key).replace("_", " ").title()))
            pdf.ln(2) 
        for item in value:
            add_json_to_pdf(pdf, "", item)
        pdf.ln(4)
    else:
        pdf.set_font("helvetica", "", 11)
        pdf.set_text_color(66, 71, 84)
        pdf.set_x(10) 
        prefix = f"{safe_text(str(key).replace('_', ' ').title())}: " if key else "- "
        final_text = safe_text(f"{prefix}{str(value)}")
        pdf.multi_cell(0, 6, txt=final_text)
        pdf.ln(2)

# Custom Worksheet Formatter
def format_worksheet(pdf, data):
    title = data.get("title", "")
    if title:
        pdf.set_font("helvetica", "B", 16)
        pdf.set_text_color(11, 28, 48)
        pdf.set_x(10)
        pdf.multi_cell(0, 8, txt=safe_text(title))
        pdf.ln(6)
        
    questions = data.get("questions", [])
    for i, q in enumerate(questions):
        q_text = str(q.get("question", ""))
        diff = str(q.get("difficulty", "Medium")).title()
        ans = str(q.get("answer_key", ""))
        
        # Print Question Number and Difficulty
        pdf.set_font("helvetica", "B", 12)
        pdf.set_text_color(0, 88, 190) # Theme Blue
        pdf.set_x(10)
        pdf.cell(15, 6, txt=f"Q{i+1}.", ln=0)
        
        pdf.set_font("helvetica", "I", 10)
        pdf.set_text_color(128, 128, 128) # Gray
        pdf.cell(0, 6, txt=f"[{diff}]", ln=1)
        
        # Print Question Text
        pdf.set_font("helvetica", "", 11)
        pdf.set_text_color(30, 41, 59) # Slate Text
        pdf.set_x(15)
        pdf.multi_cell(180, 6, txt=safe_text(q_text))
        pdf.ln(2)
        
        # Print Answer Key
        if ans:
            pdf.set_font("helvetica", "B", 10)
            pdf.set_text_color(15, 118, 110) # Teal/Emerald for Answers
            pdf.set_x(15)
            pdf.cell(18, 6, txt="Answer:", ln=0)
            
            pdf.set_font("helvetica", "", 10)
            pdf.set_text_color(71, 85, 105)
            pdf.set_x(33)
            pdf.multi_cell(165, 6, txt=safe_text(ans))
            
        pdf.ln(6)

# Custom Grading Rubric Formatter
def format_rubric(pdf, data):
    title = data.get("assignment_title", "Grading Rubric")
    score = data.get("total_score", "")
    
    if title:
        pdf.set_font("helvetica", "B", 16)
        pdf.set_text_color(11, 28, 48)
        pdf.set_x(10)
        pdf.multi_cell(0, 8, txt=safe_text(title))
        pdf.ln(2)
        
    if score:
        pdf.set_font("helvetica", "B", 11)
        pdf.set_text_color(100, 116, 139) # Slate Gray
        pdf.set_x(10)
        pdf.multi_cell(0, 6, txt=f"Total Score: {score}")
        pdf.ln(6)
        
    criteria = data.get("criteria", [])
    for i, c in enumerate(criteria):
        c_name = str(c.get("criterion_name", f"Criterion {i+1}"))
        weight = str(c.get("weight", ""))
        
        # Criterion Header
        pdf.set_font("helvetica", "B", 12)
        pdf.set_text_color(0, 88, 190) # Theme Blue
        pdf.set_x(10)
        header_text = c_name
        if weight:
            header_text += f" (Weight: {weight}%)"
        pdf.multi_cell(0, 8, txt=safe_text(header_text))
        pdf.ln(2)
        
        # Color-coded Evaluation Levels
        levels = [
            ("Excellent", c.get("excellent", ""), 15, 118, 110), # Emerald
            ("Good", c.get("good", ""), 29, 78, 216), # Blue
            ("Needs Improv.", c.get("needs_improvement", ""), 217, 119, 6), # Amber
            ("Poor", c.get("poor", ""), 225, 29, 72) # Rose Red
        ]
        
        for lvl_name, desc, r, g, b in levels:
            if desc:
                # Colored Level Label
                pdf.set_font("helvetica", "B", 10)
                pdf.set_text_color(r, g, b)
                pdf.set_x(15)
                pdf.cell(35, 6, txt=f"{lvl_name}:", ln=0)
                
                # Description Text
                pdf.set_font("helvetica", "", 10)
                pdf.set_text_color(71, 85, 105) # Slate text
                pdf.set_x(50)
                pdf.multi_cell(145, 6, txt=safe_text(desc))
        pdf.ln(5)


@router.post("/pdf")
async def export_to_pdf(request: Request):
    try:
        payload = await request.json()
        tool_name = payload.get("tool_name", "AI Document")
        content_data = payload.get("content_data", {})

        pdf = PDFReport()
        pdf.add_page()
        
        pdf.set_font("helvetica", "B", 22)
        pdf.set_text_color(11, 28, 48)
        pdf.multi_cell(0, 12, txt=safe_text(tool_name), align='C')
        pdf.ln(10)

        # The Routing Switch: Divert formatting based on the Tool Name!
        if tool_name == "Worksheet":
            format_worksheet(pdf, content_data)
        elif tool_name == "Grading Rubric":
            format_rubric(pdf, content_data)
        else:
            add_json_to_pdf(pdf, "", content_data)

        pdf_bytes = bytes(pdf.output())
        stream = io.BytesIO(pdf_bytes)

        return StreamingResponse(
            stream, 
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={tool_name.replace(' ', '_')}.pdf"}
        )
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PDF Generation failed: {str(e)}")